"""
tools/file_creator.py — File creation tool for Vera.
Creates Word (.docx), PDF, Excel (.xlsx), and PowerPoint (.pptx) files.
Vera generates the content via AI, then saves to Documents\Vera\.
"""

import os
import json
from pathlib import Path
from datetime import datetime

# Output directory — saves to user's Documents\Vera folder
OUTPUT_DIR = Path(os.path.expanduser("~")) / "Documents" / "Vera"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def _safe_filename(name: str) -> str:
    """Strip invalid filename characters."""
    import re
    name = re.sub(r'[<>:"/\\|?*]', '', name)
    return name.strip()[:60] or "vera_file"


def _timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


# ── Word (.docx) ──────────────────────────────────────────────────────────────
def create_word(filename: str = None, title: str = None, content: str = None,
                name: str = None, topic: str = None, **kwargs) -> str:
    """Create a Word document."""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return "Missing library. Run: pip install python-docx"

    filename = filename or name or topic or "document"
    title    = title or topic or filename
    content  = content or kwargs.get("text", "")

    doc  = Document()
    safe = _safe_filename(filename)
    path = OUTPUT_DIR / f"{safe}_{_timestamp()}.docx"

    # Title
    h = doc.add_heading(title, 0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Date
    p = doc.add_paragraph(datetime.now().strftime("%B %d, %Y"))
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    doc.add_paragraph()

    # Content — split by newlines, handle headings marked with ##
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            doc.add_paragraph()
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=1)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=2)
        elif line.startswith("- ") or line.startswith("• "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)

    doc.save(str(path))
    return f"Word doc saved: {path}"


# ── PDF ───────────────────────────────────────────────────────────────────────
def create_pdf(filename: str = None, title: str = None, content: str = None,
               name: str = None, topic: str = None, **kwargs) -> str:
    """Create a PDF document."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    except ImportError:
        return "Missing library. Run: pip install reportlab"

    filename = filename or name or topic or "document"
    title    = title or topic or filename
    content  = content or kwargs.get("text", "")

    safe = _safe_filename(filename)
    path = OUTPUT_DIR / f"{safe}_{_timestamp()}.pdf"

    doc    = SimpleDocTemplate(str(path), pagesize=A4,
                               leftMargin=2.5*cm, rightMargin=2.5*cm,
                               topMargin=2.5*cm, bottomMargin=2.5*cm)
    styles = getSampleStyleSheet()
    story  = []

    # Title
    title_style = ParagraphStyle("title", parent=styles["Title"],
                                 fontSize=20, spaceAfter=6)
    story.append(Paragraph(title, title_style))

    # Date
    date_style = ParagraphStyle("date", parent=styles["Normal"],
                                fontSize=10, textColor=colors.grey, spaceAfter=20)
    story.append(Paragraph(datetime.now().strftime("%B %d, %Y"), date_style))

    # Content
    h1_style = ParagraphStyle("h1", parent=styles["Heading1"], fontSize=14, spaceAfter=6)
    h2_style = ParagraphStyle("h2", parent=styles["Heading2"], fontSize=12, spaceAfter=4)

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            story.append(Spacer(1, 8))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:], h1_style))
        elif line.startswith("# "):
            story.append(Paragraph(line[2:], h2_style))
        elif line.startswith("- ") or line.startswith("• "):
            bullet = ParagraphStyle("bullet", parent=styles["Normal"],
                                    leftIndent=20, bulletIndent=10)
            story.append(Paragraph(f"• {line[2:]}", bullet))
        else:
            story.append(Paragraph(line, styles["Normal"]))

    doc.build(story)
    return f"PDF saved: {path}"


# ── Excel (.xlsx) ─────────────────────────────────────────────────────────────
def create_excel(filename: str = None, title: str = None, content: str = None,
                 name: str = None, topic: str = None, **kwargs) -> str:
    """Create an Excel spreadsheet. Content should be CSV-like rows."""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        return "Missing library. Run: pip install openpyxl"

    filename = filename or name or topic or "spreadsheet"
    title    = title or topic or filename
    content  = content or kwargs.get("text", "")

    safe = _safe_filename(filename)
    path = OUTPUT_DIR / f"{safe}_{_timestamp()}.xlsx"

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:30]

    # Header style
    header_fill = PatternFill("solid", fgColor="8B1A2E")
    header_font = Font(bold=True, color="FFFFFF", size=11)

    row_num = 1
    first_row = True

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue

        # Split by comma or tab
        if "\t" in line:
            cells = line.split("\t")
        else:
            cells = line.split(",")

        cells = [c.strip().strip('"') for c in cells]

        for col_num, value in enumerate(cells, 1):
            cell = ws.cell(row=row_num, column=col_num, value=value)
            cell.alignment = Alignment(wrap_text=True)
            if first_row:
                cell.font = header_font
                cell.fill = header_fill

        first_row = False
        row_num += 1

    # Auto-width columns
    for col in ws.columns:
        max_len = 0
        col_letter = col[0].column_letter
        for cell in col:
            if cell.value:
                max_len = max(max_len, len(str(cell.value)))
        ws.column_dimensions[col_letter].width = min(max_len + 4, 40)

    wb.save(str(path))
    return f"Excel file saved: {path}"


# ── PowerPoint (.pptx) ────────────────────────────────────────────────────────
def create_powerpoint(filename: str = None, title: str = None, content: str = None,
                      name: str = None, topic: str = None, **kwargs) -> str:
    """Create a PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return "Missing library. Run: pip install python-pptx"

    filename = filename or name or topic or "presentation"
    title    = title or topic or filename
    content  = content or kwargs.get("text", "")

    safe = _safe_filename(filename)
    path = OUTPUT_DIR / f"{safe}_{_timestamp()}.pptx"

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK   = RGBColor(0x0c, 0x0a, 0x0e)
    GARNET = RGBColor(0x8b, 0x1a, 0x2e)
    ROSE   = RGBColor(0xe8, 0x60, 0x7a)
    WHITE  = RGBColor(0xf0, 0xea, 0xf4)
    MUTED  = RGBColor(0x6b, 0x5f, 0x75)

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = DARK

    def add_title_slide(prs, title, subtitle=""):
        layout = prs.slide_layouts[6]  # blank
        slide  = prs.slides.add_slide(layout)
        set_bg(slide)

        # Title box
        txb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = ROSE

        if subtitle:
            txb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(0.8))
            tf2  = txb2.text_frame
            p2   = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = subtitle
            run2.font.size = Pt(18)
            run2.font.color.rgb = MUTED

        # Accent line
        line = slide.shapes.add_connector(1, Inches(4), Inches(4.0), Inches(9.33), Inches(4.0))
        line.line.color.rgb = GARNET
        line.line.width = Pt(2)

        return slide

    def add_content_slide(prs, slide_title, bullets):
        layout = prs.slide_layouts[6]
        slide  = prs.slides.add_slide(layout)
        set_bg(slide)

        # Title bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = GARNET
        bar.line.fill.background()

        # Slide title
        txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12), Inches(0.9))
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide_title
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Bullets
        if bullets:
            txb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(5.8))
            tf2  = txb2.text_frame
            tf2.word_wrap = True
            for i, bullet in enumerate(bullets):
                p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p2.space_before = Pt(8)
                run2 = p2.add_run()
                run2.text = f"  {bullet}"
                run2.font.size = Pt(18)
                run2.font.color.rgb = WHITE

        return slide

    # Parse content into slides
    # Format: "SLIDE: Title\n- bullet\n- bullet\nSLIDE: Next..."
    slides_data = []
    current_title   = None
    current_bullets = []

    for line in content.split("\n"):
        line = line.strip()
        if line.upper().startswith("SLIDE:"):
            if current_title is not None:
                slides_data.append((current_title, current_bullets))
            current_title   = line[6:].strip()
            current_bullets = []
        elif line.startswith("- ") or line.startswith("• "):
            current_bullets.append(line[2:])
        elif line and current_title is not None:
            current_bullets.append(line)

    if current_title is not None:
        slides_data.append((current_title, current_bullets))

    # Build slides
    if not slides_data:
        # Just make a title slide with raw content as subtitle
        add_title_slide(prs, title, content[:150])
    else:
        # First slide = title slide
        first_title, first_bullets = slides_data[0]
        subtitle = first_bullets[0] if first_bullets else ""
        add_title_slide(prs, title, subtitle)

        # Rest = content slides
        for slide_title, bullets in slides_data[1:] if len(slides_data) > 1 else slides_data:
            add_content_slide(prs, slide_title, bullets)

    prs.save(str(path))
    return f"PowerPoint saved: {path}"


# ── Open the output folder ────────────────────────────────────────────────────
def open_output_folder(**kwargs) -> str:
    """Open Vera's output folder in Explorer."""
    import subprocess
    subprocess.Popen(["explorer", str(OUTPUT_DIR)])
    return f"Opened: {OUTPUT_DIR}"